#!/usr/bin/env python3
"""
Convert a presentation HTML file to PPTX.

Usage:
    python3 tools/html_to_pptx.py final_slides/my-slides.html
    python3 tools/html_to_pptx.py final_slides/my-slides.html --output final_slides/my-slides.pptx

Output: a .pptx file saved next to the source HTML (or at --output path).

Requirements:
    pip install -r tools/requirements.txt
    playwright install chromium

How it works:
    1. Opens the HTML in a headless Chromium browser at 1920x1080
    2. Finds all .slide elements
    3. Shows each slide one at a time and screenshots it
    4. Assembles screenshots into a 16:9 PPTX using python-pptx
"""

import sys
import asyncio
import argparse
from io import BytesIO
from pathlib import Path


def parse_args():
    parser = argparse.ArgumentParser(description="Convert HTML slides to PPTX")
    parser.add_argument("html", help="Path to the HTML presentation file")
    parser.add_argument("--output", "-o", help="Output PPTX path (default: same dir as HTML)")
    return parser.parse_args()


async def capture_slides(html_path: Path) -> list[bytes]:
    try:
        from playwright.async_api import async_playwright
    except ImportError:
        print("Error: playwright not installed. Run: pip install -r tools/requirements.txt && playwright install chromium")
        sys.exit(1)

    screenshots = []
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": 1920, "height": 1080})

        url = f"file://{html_path.absolute()}"
        await page.goto(url, wait_until="networkidle")

        # Count slides
        slide_count = await page.evaluate("() => document.querySelectorAll('.slide').length")
        if slide_count == 0:
            print("Warning: no elements with class '.slide' found. Check your HTML structure.")
            await browser.close()
            return screenshots

        print(f"  Found {slide_count} slides")

        for i in range(slide_count):
            # Show only this slide — works for both fixed-position and scrollable layouts
            await page.evaluate("""
                (index) => {
                    const slides = document.querySelectorAll('.slide');
                    slides.forEach((s, i) => {
                        if (i === index) {
                            s.style.opacity = '1';
                            s.style.zIndex = '999';
                            s.style.visibility = 'visible';
                            s.style.pointerEvents = 'auto';
                            // For scroll-based layouts, scroll to the slide
                            if (getComputedStyle(s).position !== 'fixed') {
                                s.scrollIntoView({ block: 'start' });
                            }
                        } else {
                            s.style.opacity = '0';
                            s.style.zIndex = '0';
                            s.style.visibility = 'hidden';
                            s.style.pointerEvents = 'none';
                        }
                    });
                    // Hide nav dots so they don't appear in the screenshot
                    document.querySelectorAll('.nav, .dots, .navigation, .slide-nav').forEach(n => n.style.display = 'none');
                }
            """, i)

            await page.wait_for_timeout(300)  # Let transitions settle
            screenshot = await page.screenshot(full_page=False)
            screenshots.append(screenshot)
            print(f"  Captured slide {i + 1}/{slide_count}")

        await browser.close()

    return screenshots


def build_pptx(screenshots: list[bytes], output_path: Path):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("Error: python-pptx not installed. Run: pip install -r tools/requirements.txt")
        sys.exit(1)

    prs = Presentation()
    # 16:9 at standard PPTX dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Completely blank

    for i, screenshot_bytes in enumerate(screenshots):
        slide = prs.slides.add_slide(blank_layout)
        img_stream = BytesIO(screenshot_bytes)
        slide.shapes.add_picture(
            img_stream,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        print(f"  Added slide {i + 1}/{len(screenshots)} to PPTX")

    prs.save(output_path)


async def main():
    args = parse_args()

    html_path = Path(args.html).resolve()
    if not html_path.exists():
        print(f"Error: file not found: {html_path}")
        sys.exit(1)

    if args.output:
        output_path = Path(args.output).resolve()
    else:
        output_path = html_path.with_suffix(".pptx")

    print(f"Source:  {html_path}")
    print(f"Output:  {output_path}")
    print()

    print("Capturing slides...")
    screenshots = await capture_slides(html_path)

    if not screenshots:
        print("No slides captured. Exiting.")
        sys.exit(1)

    print(f"Captured {len(screenshots)} slides")
    print()

    print("Building PPTX...")
    build_pptx(screenshots, output_path)

    size_kb = output_path.stat().st_size // 1024
    print()
    print(f"Done! Saved to: {output_path} ({size_kb} KB)")


if __name__ == "__main__":
    asyncio.run(main())
